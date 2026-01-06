"""
PowerPoint Slide Extractor

This script extracts slides from a PowerPoint presentation as PNG images
and saves speaker notes as text files.

Usage:
    python extract_slides.py <path_to_presentation.pptx>
"""

import argparse
import logging
import os
import subprocess
import sys
from pathlib import Path
from typing import Optional

import comtypes.client
from pptx import Presentation


# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('error.log'),
        logging.StreamHandler(sys.stdout)
    ]
)
logger = logging.getLogger(__name__)


# Video configuration (adjust for quality vs speed tradeoff)
VIDEO_CONFIG = {
    'fps': 24,             # Frames per second (24 is standard for video, even with static images)
    'codec': 'libx264',    # Video codec
    'audio_codec': 'aac',  # Audio codec
    'preset': 'ultrafast', # Encoding speed: ultrafast, superfast, veryfast, faster, fast, medium, slow, slower, veryslow
    'bitrate': '1000k',    # Video bitrate (higher = better quality, larger file)
}


def setup_output_folder(pptx_path: Path) -> Path:
    """
    Create output folder with the same name as the presentation.
    
    Args:
        pptx_path: Path to the PowerPoint file
        
    Returns:
        Path to the output folder
    """
    output_folder = pptx_path.parent / pptx_path.stem
    
    # Create folder (overwrite contents if exists)
    output_folder.mkdir(parents=True, exist_ok=True)
    logger.info(f"Output folder: {output_folder}")
    
    return output_folder


def export_slides_as_png(pptx_path: Path, output_folder: Path) -> int:
    """
    Export all slides from PowerPoint presentation as PNG images using COM automation.
    
    Args:
        pptx_path: Path to the PowerPoint file
        output_folder: Path to the output folder
        
    Returns:
        Number of slides exported
    """
    powerpoint = None
    presentation = None
    slide_count = 0
    
    try:
        # Convert to absolute path for COM
        abs_pptx_path = str(pptx_path.resolve())
        
        logger.info("Starting PowerPoint application...")
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1  # Make PowerPoint visible
        
        logger.info(f"Opening presentation: {abs_pptx_path}")
        presentation = powerpoint.Presentations.Open(abs_pptx_path, WithWindow=False)
        
        slide_count = presentation.Slides.Count
        logger.info(f"Found {slide_count} slides")
        
        # Export each slide as PNG (skip hidden slides)
        visible_slide_num = 0
        for i in range(1, slide_count + 1):
            try:
                slide = presentation.Slides(i)
                
                # Check if slide is hidden
                try:
                    if slide.SlideShowTransition.Hidden:
                        logger.info(f"Skipping hidden slide {i}")
                        continue
                except:
                    pass  # If we can't check, assume it's visible
                
                visible_slide_num += 1
                output_file = output_folder / f"slide_{visible_slide_num:02d}.png"
                
                logger.info(f"Exporting slide {i} as slide {visible_slide_num}/{slide_count} to {output_file.name}")
                slide.Export(str(output_file.resolve()), "PNG")
                
            except Exception as e:
                logger.error(f"Failed to export slide {i}: {e}")
        
        logger.info(f"Slide export completed - exported {visible_slide_num} visible slides")
        
    except Exception as e:
        logger.error(f"Error during PowerPoint automation: {e}")
        raise
        
    finally:
        # Clean up COM objects
        try:
            if presentation:
                presentation.Close()
                logger.info("Presentation closed")
        except Exception as e:
            logger.error(f"Error closing presentation: {e}")
            
        try:
            if powerpoint:
                powerpoint.Quit()
                logger.info("PowerPoint application closed")
        except Exception as e:
            logger.error(f"Error quitting PowerPoint: {e}")
    
    return visible_slide_num


def extract_speaker_notes(pptx_path: Path, output_folder: Path, slide_count: int) -> None:
    """
    Extract speaker notes from each visible slide and save as text files.
    Skips hidden slides.
    
    Args:
        pptx_path: Path to the PowerPoint file
        output_folder: Path to the output folder
        slide_count: Expected number of visible slides
    """
    try:
        logger.info("Extracting speaker notes...")
        prs = Presentation(str(pptx_path))
        
        visible_slide_num = 0
        for i, slide in enumerate(prs.slides, start=1):
            try:
                # Check if slide is hidden using the underlying XML
                try:
                    if hasattr(slide, '_element') and hasattr(slide._element, 'show'):
                        if slide._element.show == 0:
                            logger.info(f"Skipping hidden slide {i} for notes extraction")
                            continue
                except:
                    pass  # If we can't check, assume it's visible
                
                visible_slide_num += 1
                
                # Get speaker notes
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text if notes_slide else ""
                
                # Save to text file (only if content has changed)
                output_file = output_folder / f"text_{visible_slide_num:02d}.txt"
                
                # Check if file exists and content is unchanged
                should_write = True
                if output_file.exists():
                    existing_text = output_file.read_text(encoding='utf-8')
                    if existing_text == notes_text:
                        should_write = False
                        logger.info(f"Text {visible_slide_num} unchanged, preserving timestamp for caching")
                
                if should_write:
                    output_file.write_text(notes_text, encoding='utf-8')
                    if notes_text.strip():
                        logger.info(f"Extracted notes for slide {i} as text_{visible_slide_num:02d}.txt")
                    else:
                        logger.info(f"No notes found for slide {i} (text_{visible_slide_num:02d}.txt)")
                    
            except Exception as e:
                logger.error(f"Failed to extract notes for slide {i}: {e}")
        
        logger.info(f"Speaker notes extraction completed - processed {visible_slide_num} visible slides")
        
    except Exception as e:
        logger.error(f"Error reading presentation for notes: {e}")
        raise


def generate_audio_from_notes(output_folder: Path, slide_count: int) -> None:
    """
    Generate audio files from speaker notes text files using csm-voice tool.
    Caches audio files to avoid regeneration if text hasn't changed.
    
    Args:
        output_folder: Path to the output folder containing text files
        slide_count: Number of slides in the presentation
    """
    try:
        logger.info("Generating audio from speaker notes...")
        
        for i in range(1, slide_count + 1):
            text_file = f"text_{i:02d}.txt"
            audio_file = f"audio_{i:02d}.wav"
            
            text_path = output_folder / text_file
            audio_path = output_folder / audio_file
            
            try:
                # Check if audio already exists and is newer than text file
                if audio_path.exists():
                    audio_mtime = audio_path.stat().st_mtime
                    text_mtime = text_path.stat().st_mtime
                    
                    if audio_mtime > text_mtime:
                        logger.info(f"Audio {i}/{slide_count} already exists and is up-to-date, reusing...")
                        continue
                
                logger.info(f"Processing {text_file} -> {audio_file}")
                
                # Set up environment with UTF-8 encoding for csm-voice
                env = os.environ.copy()
                env['PYTHONIOENCODING'] = 'utf-8'
                
                # Run csm-voice from the output folder
                result = subprocess.run(
                    [r"D:\Dev\lesterthomas\csm-lester-voice\.venv\Scripts\csm-voice.exe", "-f", text_file, "-o", audio_file],
                    cwd=str(output_folder),
                    capture_output=True,
                    text=True,
                    check=True,
                    env=env
                )
                
                logger.info(f"Successfully generated {audio_file}")
                
            except subprocess.CalledProcessError as e:
                logger.error(f"Failed to generate audio for {text_file}: {e}")
                if e.stderr:
                    logger.error(f"Error output: {e.stderr}")
            except FileNotFoundError:
                logger.error("csm-voice tool not found. Please ensure it is installed and in your PATH.")
                break
            except Exception as e:
                logger.error(f"Unexpected error processing {text_file}: {e}")
        
        logger.info("Audio generation completed")
        
    except Exception as e:
        logger.error(f"Error during audio generation: {e}")
        raise


def create_individual_clips(output_folder: Path, slide_count: int) -> list:
    """
    Create individual video clips for each slide (image + audio).
    Caches clips to avoid regenerating if files haven't changed.
    
    Args:
        output_folder: Path to the output folder containing slides and audio
        slide_count: Number of slides in the presentation
        
    Returns:
        List of paths to created video clips
    """
    try:
        from moviepy import ImageClip, AudioFileClip, CompositeAudioClip
        
        logger.info("Creating individual video clips...")
        
        clips_folder = output_folder / "clips"
        clips_folder.mkdir(exist_ok=True)
        
        clip_paths = []
        
        for i in range(1, slide_count + 1):
            try:
                image_file = output_folder / f"slide_{i:02d}.png"
                audio_file = output_folder / f"audio_{i:02d}.wav"
                clip_file = clips_folder / f"clip_{i:02d}.mp4"
                
                # Check if clip already exists and is newer than source files
                if clip_file.exists():
                    clip_mtime = clip_file.stat().st_mtime
                    image_mtime = image_file.stat().st_mtime
                    audio_mtime = audio_file.stat().st_mtime
                    
                    if clip_mtime > image_mtime and clip_mtime > audio_mtime:
                        logger.info(f"Clip {i}/{slide_count} already exists and is up-to-date, reusing...")
                        clip_paths.append(clip_file)
                        continue
                
                logger.info(f"Creating video clip {i}/{slide_count}...")
                
                audio = AudioFileClip(str(audio_file))
                # Add 1 second pause at the start of each slide
                # Set audio to start 1 second into the clip
                audio_delayed = audio.with_start(1.0)
                # Create composite with just the delayed audio
                composite_audio = CompositeAudioClip([audio_delayed])
                # Image duration is audio duration + 1 second pause
                image = ImageClip(str(image_file)).with_duration(audio.duration + 1.0)
                video = image.with_audio(composite_audio)
                
                video.write_videofile(
                    str(clip_file),
                    fps=VIDEO_CONFIG['fps'],
                    codec=VIDEO_CONFIG['codec'],
                    audio_codec=VIDEO_CONFIG['audio_codec'],
                    preset=VIDEO_CONFIG['preset'],
                    bitrate=VIDEO_CONFIG['bitrate'],
                    logger=None  # Suppress MoviePy progress bars
                )
                
                video.close()
                audio.close()
                
                clip_paths.append(clip_file)
                logger.info(f"Successfully created clip {i}")
                
            except Exception as e:
                logger.error(f"Failed to create clip {i}: {e}")
                raise
        
        logger.info("Individual clip creation completed")
        return clip_paths
        
    except Exception as e:
        logger.error(f"Error during clip creation: {e}")
        raise


def concatenate_clips(clip_paths: list, output_file: Path) -> None:
    """
    Concatenate all video clips into a single final video using FFmpeg.
    
    Args:
        clip_paths: List of paths to individual video clips
        output_file: Path where the final video should be saved
    """
    try:
        import imageio_ffmpeg
        
        logger.info(f"Concatenating {len(clip_paths)} clips into final video...")
        
        # Create a temporary file list for FFmpeg concat
        concat_file = output_file.parent / "concat_list.txt"
        with open(concat_file, 'w') as f:
            for clip_path in clip_paths:
                # Use forward slashes and escape special characters for FFmpeg
                safe_path = str(clip_path.absolute()).replace('\\', '/')
                f.write(f"file '{safe_path}'\n")
        
        # Get FFmpeg executable from imageio_ffmpeg
        ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()
        
        # Use FFmpeg directly for concatenation
        result = subprocess.run(
            [
                ffmpeg_exe, '-f', 'concat', '-safe', '0', '-i', str(concat_file),
                '-c', 'copy', str(output_file), '-y'
            ],
            capture_output=True,
            text=True
        )
        
        # Clean up temp file
        concat_file.unlink()
        
        if result.returncode != 0:
            logger.error(f"FFmpeg error: {result.stderr}")
            raise RuntimeError(f"FFmpeg concatenation failed: {result.stderr}")
        
        logger.info(f"Final video created: {output_file}")
        
    except Exception as e:
        logger.error(f"Error during video concatenation: {e}")
        raise


def create_presentation_video(output_folder: Path, slide_count: int) -> None:
    """
    Create final presentation video from slides and audio files.
    
    Args:
        output_folder: Path to the output folder
        slide_count: Number of slides in the presentation
    """
    try:
        logger.info("Starting video creation process...")
        
        # Step 1: Create individual clips (with caching)
        clip_paths = create_individual_clips(output_folder, slide_count)
        
        # Step 2: Concatenate into final video
        final_video_path = output_folder / f"{output_folder.name}_video.mp4"
        concatenate_clips(clip_paths, final_video_path)
        
        logger.info("Video creation completed successfully!")
        
    except Exception as e:
        logger.error(f"Failed to create presentation video: {e}")
        raise


def main():
    """Main entry point for the script."""
    parser = argparse.ArgumentParser(
        description="Extract slides and speaker notes from PowerPoint presentations"
    )
    parser.add_argument(
        "presentation",
        type=str,
        help="Path to the PowerPoint presentation file (.pptx)"
    )
    
    args = parser.parse_args()
    
    # Validate input file
    pptx_path = Path(args.presentation)
    if not pptx_path.exists():
        logger.error(f"File not found: {pptx_path}")
        sys.exit(1)
    
    if not pptx_path.suffix.lower() in ['.pptx', '.ppt']:
        logger.error(f"Invalid file type: {pptx_path.suffix}. Expected .pptx or .ppt")
        sys.exit(1)
    
    logger.info(f"Processing presentation: {pptx_path}")
    
    try:
        # Set up output folder
        output_folder = setup_output_folder(pptx_path)
        
        # Export slides as PNG
        slide_count = export_slides_as_png(pptx_path, output_folder)
        
        # Extract speaker notes
        extract_speaker_notes(pptx_path, output_folder, slide_count)
        
        # Generate audio from speaker notes
        generate_audio_from_notes(output_folder, slide_count)
        
        # Create presentation video
        create_presentation_video(output_folder, slide_count)
        
        logger.info(f"Successfully processed {slide_count} slides")
        logger.info(f"Output saved to: {output_folder}")
        
    except Exception as e:
        logger.error(f"Failed to process presentation: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
